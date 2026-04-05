from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Pt


PROJECT_ROOT = Path(__file__).resolve().parent
FINAL_ROOT = PROJECT_ROOT / "final_assets"
FIG_ROOT = FINAL_ROOT / "figures"
PPT_ROOT = FINAL_ROOT / "ppt"
PPT_PATH = PPT_ROOT / "水稻叶病害分类答辩.pptx"


FONT_NAME = "Microsoft YaHei"
COLOR_DARK = RGBColor(28, 68, 52)
COLOR_ACCENT = RGBColor(197, 152, 62)
COLOR_LIGHT = RGBColor(247, 245, 239)
COLOR_TEXT = RGBColor(34, 40, 49)
COLOR_MUTED = RGBColor(83, 91, 104)
COLOR_WHITE = RGBColor(255, 255, 255)


def add_textbox(slide, left, top, width, height, text="", font_size=20, bold=False, color=COLOR_TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT_NAME
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, bullets, font_size=18, color=COLOR_TEXT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, item in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"• {item}"
        paragraph.level = 0
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.font.name = FONT_NAME
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = color
        paragraph.space_after = Pt(8)
    return box


def add_title_band(slide, title, subtitle=None):
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Cm(1.6)).fill.solid()
    band = slide.shapes[-1]
    band.fill.fore_color.rgb = COLOR_DARK
    band.line.color.rgb = COLOR_DARK
    add_textbox(slide, Cm(0.8), Cm(0.18), Cm(20), Cm(0.9), title, font_size=26, bold=True, color=COLOR_WHITE)
    if subtitle:
        add_textbox(slide, Cm(21), Cm(0.28), Cm(11), Cm(0.7), subtitle, font_size=11, color=RGBColor(226, 232, 240), align=PP_ALIGN.RIGHT)


def add_background(slide, image_path: Path):
    slide.shapes.add_picture(str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    overlay = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(17, 24, 39)
    overlay.fill.transparency = 0.35
    overlay.line.color.rgb = RGBColor(17, 24, 39)


def add_card(slide, left, top, width, height, title, value=None):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_LIGHT
    card.line.color.rgb = COLOR_ACCENT
    add_textbox(slide, left + Cm(0.4), top + Cm(0.25), width - Cm(0.8), Cm(0.8), title, font_size=13, bold=True, color=COLOR_MUTED)
    if value is not None:
        add_textbox(slide, left + Cm(0.4), top + Cm(0.95), width - Cm(0.8), Cm(0.9), value, font_size=22, bold=True, color=COLOR_DARK)


def add_image(slide, image_path: Path, left, top, width=None, height=None):
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def add_table(slide, left, top, width, height, headers, rows, font_size=12):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
    table = table_shape.table
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_DARK
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = FONT_NAME
                run.font.bold = True
                run.font.size = Pt(font_size)
                run.font.color.rgb = COLOR_WHITE
            paragraph.alignment = PP_ALIGN.CENTER
    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            fill_color = COLOR_LIGHT if row_idx % 2 == 1 else RGBColor(236, 240, 232)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill_color
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = FONT_NAME
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = COLOR_TEXT
                paragraph.alignment = PP_ALIGN.CENTER


def slide_cover():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, FIG_ROOT / "official_best_gradcam_gallery.png")
    ribbon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(1.0), Cm(2.0), Cm(16.5), Cm(6.5))
    ribbon.fill.solid()
    ribbon.fill.fore_color.rgb = RGBColor(20, 32, 28)
    ribbon.fill.transparency = 0.08
    ribbon.line.color.rgb = COLOR_ACCENT
    add_textbox(slide, Cm(1.8), Cm(2.8), Cm(15), Cm(1.8), "水稻叶病害图像分类课程项目答辩", font_size=28, bold=True, color=COLOR_WHITE)
    add_textbox(slide, Cm(1.9), Cm(4.9), Cm(13), Cm(3.0), "课程：________\n姓名：________\n学号：________\n日期：2026-04-03", font_size=16, color=RGBColor(232, 240, 236))
    footer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, prs.slide_height - Cm(1.2), prs.slide_width, Cm(1.2))
    footer.fill.solid()
    footer.fill.fore_color.rgb = COLOR_ACCENT
    footer.line.color.rgb = COLOR_ACCENT
    add_textbox(slide, Cm(0.8), prs.slide_height - Cm(0.95), Cm(28), Cm(0.6), "关键词：双轨评测 / 迁移学习 / 超参数调优 / 稳定性验证", font_size=14, bold=True, color=COLOR_DARK)


def slide_task():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_band(slide, "研究背景与任务定义", "任务页")
    add_bullets(slide, Cm(0.9), Cm(2.1), Cm(12.2), Cm(7.0), [
        "任务目标：根据叶片图像识别 4 类水稻叶病害。",
        "课程目标：构建高准确率、可复现、可解释的完整课程项目。",
        "研究对象：Bacterialblight、Blast、Brownspot、Tungro。",
        "项目交付：训练代码、调参结果、可视化图表、中文报告与答辩材料。",
    ], font_size=18)
    add_image(slide, FIG_ROOT / "official_class_distribution.png", Cm(14.2), Cm(2.0), width=Cm(18))
    add_card(slide, Cm(14.4), Cm(8.2), Cm(5.0), Cm(2.0), "类别数", "4 类")
    add_card(slide, Cm(20.0), Cm(8.2), Cm(5.5), Cm(2.0), "任务类型", "图像四分类")
    add_card(slide, Cm(26.0), Cm(8.2), Cm(5.2), Cm(2.0), "项目定位", "课程答辩")


def slide_dataset():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_band(slide, "数据集与关键问题", "数据审计页")
    add_bullets(slide, Cm(0.9), Cm(2.0), Cm(12.0), Cm(7.5), [
        "原始数据总量：5932 张图像，4 个病害类别。",
        "官方划分：train 4700 / validation 616 / test 616。",
        "哈希审计发现：224 组跨集合重复图，共涉及 448 张图片。",
        "结论：如果只看官方划分，结果会偏乐观，需要补充更严谨的 clean 视图。",
    ], font_size=18)
    add_image(slide, FIG_ROOT / "duplicate_overview.png", Cm(14.0), Cm(2.0), width=Cm(17.5))
    add_card(slide, Cm(14.3), Cm(8.2), Cm(5.2), Cm(2.0), "总图片数", "5932")
    add_card(slide, Cm(20.2), Cm(8.2), Cm(5.2), Cm(2.0), "唯一哈希数", "4794")
    add_card(slide, Cm(26.1), Cm(8.2), Cm(5.2), Cm(2.0), "跨集合重复组", "224")


def slide_dual_view():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_band(slide, "双轨评测设计", "方法设计页")
    add_bullets(slide, Cm(0.9), Cm(2.1), Cm(10.5), Cm(3.6), [
        "official：保留原始官方划分，用于展示高分结果。",
        "clean：按图像内容哈希去重后重新分层划分，用于评估真实泛化能力。",
        "双轨设计的目标是兼顾课程成绩展示与实验严谨性。",
    ], font_size=18)
    add_image(slide, FIG_ROOT / "official_class_distribution.png", Cm(1.0), Cm(5.4), width=Cm(14.5))
    add_image(slide, FIG_ROOT / "clean_class_distribution.png", Cm(16.5), Cm(5.4), width=Cm(14.5))
    add_card(slide, Cm(12.0), Cm(2.3), Cm(9.0), Cm(2.0), "official", "展示高分结果")
    add_card(slide, Cm(22.0), Cm(2.3), Cm(9.0), Cm(2.0), "clean", "验证泛化能力")


def slide_pipeline():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_band(slide, "技术路线总览", "流程页")
    steps = [
        ("数据审计", "统计样本分布、检查重复图"),
        ("模型训练", "CustomCNN / ResNet18 / EfficientNet-B0"),
        ("超参数调优", "official 视图下对 EfficientNet-B0 做 Optuna 搜索"),
        ("测试评估", "输出 Accuracy、Macro F1、混淆矩阵、错误样本、Grad-CAM"),
    ]
    x_positions = [Cm(1.0), Cm(8.5), Cm(16.0), Cm(23.5)]
    for (title, desc), left in zip(steps, x_positions):
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, Cm(3.0), Cm(6.0), Cm(3.6))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_LIGHT
        card.line.color.rgb = COLOR_ACCENT
        add_textbox(slide, left + Cm(0.35), Cm(3.35), Cm(5.3), Cm(0.8), title, font_size=18, bold=True, color=COLOR_DARK, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Cm(0.35), Cm(4.25), Cm(5.3), Cm(1.8), desc, font_size=12, color=COLOR_MUTED, align=PP_ALIGN.CENTER)
    for start in [Cm(7.0), Cm(14.5), Cm(22.0)]:
        arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, start, Cm(4.0), Cm(1.0), Cm(1.1))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = COLOR_ACCENT
        arrow.line.color.rgb = COLOR_ACCENT
    add_bullets(slide, Cm(3.0), Cm(7.7), Cm(26.0), Cm(2.2), [
        "本项目不是只展示最终分数，而是完整覆盖了数据、方法、结果和可解释性。",
        "答辩时建议用这页作为全局导航页，帮助老师快速建立整体印象。",
    ], font_size=16)


def slide_model_training():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_band(slide, "模型与训练策略", "建模页")
    add_bullets(slide, Cm(0.9), Cm(1.9), Cm(10.5), Cm(3.5), [
        "CustomCNN：课程基础模型，突出自主搭建网络结构。",
        "ResNet18：轻量级迁移学习模型，clean 视图下表现最亮眼。",
        "EfficientNet-B0：调参主线模型，负责展示完整方法论。",
        "训练策略：两阶段微调 + 早停 + 学习率调度 + GPU 加速。",
    ], font_size=17)
    headers = ["模型", "视图", "batch", "epoch", "备注"]
    rows = [
        ["CustomCNN", "official/clean", "64", "20", "直接训练"],
        ["ResNet18", "official/clean", "48", "3 + 10", "两阶段微调"],
        ["EfficientNet-B0", "official/clean", "16", "2 + 5", "调参后配置"],
    ]
    add_table(slide, Cm(12.2), Cm(2.0), Cm(19.0), Cm(4.8), headers, rows, font_size=11)
    add_card(slide, Cm(12.5), Cm(7.5), Cm(5.5), Cm(2.0), "输入尺寸", "224 × 224")
    add_card(slide, Cm(18.6), Cm(7.5), Cm(5.8), Cm(2.0), "硬件环境", "RTX 4060")
    add_card(slide, Cm(24.9), Cm(7.5), Cm(6.0), Cm(2.0), "训练风格", "迁移学习主线")


def slide_tuning():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_band(slide, "EfficientNet-B0 调参结果", "调参页")
    add_image(slide, FIG_ROOT / "tuning_history.png", Cm(0.9), Cm(2.0), width=Cm(14.5))
    add_image(slide, FIG_ROOT / "best_params.png", Cm(16.4), Cm(2.0), width=Cm(14.5))
    add_bullets(slide, Cm(0.9), Cm(8.0), Cm(14.0), Cm(2.4), [
        "调参对象：official + EfficientNet-B0",
        "总 trial 数：12",
        "完成 7 次，剪枝 5 次",
    ], font_size=16)
    add_bullets(slide, Cm(16.5), Cm(8.0), Cm(14.0), Cm(2.8), [
        "最优参数：AdamW / lr=0.000271 / wd=1.296e-05",
        "dropout=0.3727 / label_smoothing=0.1128",
        "freeze=2 / finetune=5 / batch=16",
    ], font_size=16)


def slide_comparison():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_band(slide, "正式结果总对比", "结果页")
    add_image(slide, FIG_ROOT / "model_comparison_accuracy.png", Cm(0.9), Cm(1.8), width=Cm(15.0))
    add_image(slide, FIG_ROOT / "model_comparison_macro_f1.png", Cm(16.1), Cm(1.8), width=Cm(15.0))
    headers = ["视图", "模型", "Accuracy", "Macro F1"]
    rows = [
        ["official", "CustomCNN", "99.51%", "99.52%"],
        ["official", "ResNet18", "99.51%", "99.48%"],
        ["official", "EfficientNet-B0", "99.84%", "99.83%"],
        ["clean", "CustomCNN", "99.37%", "99.33%"],
        ["clean", "ResNet18", "100.00%", "100.00%"],
        ["clean", "EfficientNet-B0", "99.79%", "99.80%"],
    ]
    add_table(slide, Cm(3.0), Cm(7.8), Cm(25.0), Cm(3.0), headers, rows, font_size=11)


def slide_official_main():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_band(slide, "主线结果：official EfficientNet-B0", "official 主结果")
    add_image(slide, FIG_ROOT / "official_best_confusion_matrix.png", Cm(0.9), Cm(2.0), width=Cm(14.5))
    add_image(slide, FIG_ROOT / "official_best_class_metrics.png", Cm(16.3), Cm(2.0), width=Cm(14.6))
    add_card(slide, Cm(1.0), Cm(8.2), Cm(4.8), Cm(1.8), "Accuracy", "99.84%")
    add_card(slide, Cm(6.1), Cm(8.2), Cm(4.8), Cm(1.8), "Macro F1", "99.83%")
    add_card(slide, Cm(11.2), Cm(8.2), Cm(4.8), Cm(1.8), "验证最高", "100.00%")
    add_bullets(slide, Cm(17.0), Cm(8.1), Cm(13.2), Cm(2.3), [
        "official 口径下的最佳主线结果。",
        "仅 1 张 Bacterialblight 被误判为 Blast。",
        "Brownspot 与 Tungro 测试集全部识别正确。",
    ], font_size=15)


def slide_clean_highlight():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_band(slide, "亮点结果：clean ResNet18", "答辩亮点")
    add_image(slide, FIG_ROOT / "clean_resnet18_confusion_matrix.png", Cm(1.0), Cm(2.0), width=Cm(15.0))
    add_image(slide, FIG_ROOT / "clean_resnet18_class_metrics.png", Cm(17.0), Cm(2.2), width=Cm(13.8))
    add_card(slide, Cm(1.1), Cm(8.2), Cm(4.8), Cm(1.8), "Accuracy", "100.00%")
    add_card(slide, Cm(6.2), Cm(8.2), Cm(4.8), Cm(1.8), "Macro F1", "100.00%")
    add_card(slide, Cm(11.3), Cm(8.2), Cm(4.8), Cm(1.8), "定位", "clean 最高分")
    add_bullets(slide, Cm(17.2), Cm(8.1), Cm(13.0), Cm(2.2), [
        "当前 clean 视图下得分最高的模型。",
        "模型容量适中，和当前数据规模匹配更好。",
        "适合作为答辩重点亮点进行讲解。",
    ], font_size=15)


def slide_stability_cam():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_band(slide, "稳定性与可解释性", "可信度页")
    add_image(slide, FIG_ROOT / "resnet18_seed_consistency.png", Cm(0.9), Cm(2.0), width=Cm(14.5))
    add_image(slide, FIG_ROOT / "clean_resnet18_gradcam_gallery.png", Cm(16.2), Cm(2.0), width=Cm(14.7))
    add_bullets(slide, Cm(1.0), Cm(8.1), Cm(14.2), Cm(2.6), [
        "三次 seed 结果：100.00% / 99.79% / 100.00%。",
        "平均 Accuracy：99.93%，标准差：0.10%。",
        "结论：高分不是单次偶然，而是稳定高分。",
    ], font_size=15)
    add_bullets(slide, Cm(16.5), Cm(8.1), Cm(14.0), Cm(2.4), [
        "Grad-CAM 热区主要集中在叶片病斑区域。",
        "说明模型关注的是病害本身，而不是背景噪声。",
    ], font_size=15)


def slide_conclusion():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_band(slide, "结论与答辩预设问答", "收尾页")
    add_card(slide, Cm(0.9), Cm(2.0), Cm(9.7), Cm(2.4), "结论 1", "迁移学习整体优于自建 CNN")
    add_card(slide, Cm(11.2), Cm(2.0), Cm(9.7), Cm(2.4), "结论 2", "official 更高，但存在重复图带来的乐观偏差")
    add_card(slide, Cm(21.5), Cm(2.0), Cm(9.7), Cm(2.4), "结论 3", "clean 更可信，ResNet18 在 clean 上最亮眼")
    add_bullets(slide, Cm(1.0), Cm(5.0), Cm(14.5), Cm(4.5), [
        "为什么 official 分数通常更高？",
        "为什么 ResNet18 在 clean 上优于 EfficientNet-B0？",
        "100% 是否可信？",
        "为什么仍保留 EfficientNet 作为调参主线？",
    ], font_size=17)
    add_bullets(slide, Cm(16.5), Cm(5.0), Cm(14.2), Cm(4.5), [
        "答辩建议：先讲数据问题，再讲模型路线，最后讲亮点结果。",
        "主线结果：official EfficientNet 99.84%。",
        "亮点结果：clean ResNet18 平均 99.93%。",
    ], font_size=16)
    footer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, prs.slide_height - Cm(1.0), prs.slide_width, Cm(1.0))
    footer.fill.solid()
    footer.fill.fore_color.rgb = COLOR_DARK
    footer.line.color.rgb = COLOR_DARK
    add_textbox(slide, Cm(0.8), prs.slide_height - Cm(0.82), Cm(28), Cm(0.5), "谢谢各位老师，请批评指正。", font_size=16, bold=True, color=COLOR_WHITE)


def build():
    PPT_ROOT.mkdir(parents=True, exist_ok=True)
    slide_cover()
    slide_task()
    slide_dataset()
    slide_dual_view()
    slide_pipeline()
    slide_model_training()
    slide_tuning()
    slide_comparison()
    slide_official_main()
    slide_clean_highlight()
    slide_stability_cam()
    slide_conclusion()
    prs.save(str(PPT_PATH))
    return PPT_PATH


if __name__ == "__main__":
    prs = Presentation()
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)
    output = build()
    print(output)
